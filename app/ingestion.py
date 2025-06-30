"""
all the imports 
"""
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF

"""
Functions to access a input PDF,PPT,DOCX document type

"""

def extract_from_pdf(path):
    text = ""
    doc = fitz.open(path)
    for page in doc:
        text += page.get_text()
    return text

def extract_from_docx(path):
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_from_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
